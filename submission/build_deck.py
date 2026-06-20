"""Generates the ParkSense pitch deck (PPTX) with simple language and real screenshots."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "snapshots")
OUT = os.path.join(HERE, "ParkSense-Pitch-Deck.pptx")

# Palette (matches the product)
NAVY = RGBColor(0x0B, 0x10, 0x20)
PANEL = RGBColor(0x12, 0x19, 0x2E)
WHITE = RGBColor(0xF2, 0xF4, 0xF8)
MUTE = RGBColor(0x9A, 0xA6, 0xC0)
YELLOW = RGBColor(0xFF, 0xC2, 0x00)
BLUE = RGBColor(0x2D, 0x87, 0xF0)
RED = RGBColor(0xFF, 0x5A, 0x5A)
GREEN = RGBColor(0x46, 0xD3, 0x7E)
FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide(bg=NAVY):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0):
    shp = s.shapes.add_shape(1, x, y, w, h)  # rectangle
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = FONT
    return tb


def kicker(s, label, color=YELLOW):
    bar = box(s, Inches(0.7), Inches(0.62), Inches(0.16), Inches(0.42), fill=color)
    text(s, Inches(0.98), Inches(0.6), Inches(9), Inches(0.5),
         [[(label.upper(), 14, MUTE, True)]], anchor=MSO_ANCHOR.MIDDLE)


def title(s, t, y=Inches(1.05), size=34, color=WHITE, w=Inches(11.9)):
    text(s, Inches(0.7), y, w, Inches(1.0), [[(t, size, color, True)]])


def footer(s, n):
    text(s, Inches(0.7), Inches(7.02), Inches(8), Inches(0.35),
         [[("ParkSense  ·  Gridlock Hackathon 2.0", 9, MUTE, False)]])
    text(s, Inches(12.0), Inches(7.02), Inches(0.9), Inches(0.35),
         [[(str(n), 9, MUTE, False)]], align=PP_ALIGN.RIGHT)


def pic_framed(s, path, x, y, w):
    """Add picture scaled to width w, with a thin frame; returns placed height."""
    from PIL import Image
    iw, ih = Image.open(path).size
    h = Emu(int(w * ih / iw))
    box(s, x - Pt(1), y - Pt(1), w + Pt(2), h + Pt(2), fill=None, line=BLUE, line_w=1.2)
    s.shapes.add_picture(path, x, y, width=w)
    return h


def bullet(label, body):
    return [(label + "  ", 17, YELLOW, True), (body, 17, WHITE, False)]


# ---------------------------------------------------------------- Slide 1: Title
s = slide()
box(s, 0, 0, W, Inches(0.16), fill=YELLOW)
text(s, Inches(0.7), Inches(2.15), Inches(12), Inches(1.0),
     [[("ParkSense", 60, WHITE, True)]])
text(s, Inches(0.7), Inches(3.25), Inches(12), Inches(0.7),
     [[("AI parking enforcement intelligence for Bengaluru", 24, YELLOW, False)]])
text(s, Inches(0.74), Inches(4.15), Inches(12), Inches(0.6),
     [[("From reactive patrols to data-ranked enforcement.", 18, MUTE, False)]])
box(s, Inches(0.7), Inches(5.15), Inches(5.6), Pt(1.4), fill=PANEL, line=None)
text(s, Inches(0.7), Inches(5.35), Inches(12), Inches(1.2), [
    [("Theme:  Poor visibility on parking-induced congestion", 14, MUTE, False)],
    [("Gridlock Hackathon 2.0  ·  Round 2", 14, MUTE, False)],
])

# ---------------------------------------------------------------- Slide 2: Problem
s = slide()
kicker(s, "The problem")
title(s, "Parking chaos is choking our roads")
text(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.5), [
    bullet("Cars block the road.",
           "Illegal and double parking near markets, metro stations and event"),
    [("", 17, WHITE, False), ("venues squeezes traffic into fewer lanes.", 17, WHITE, False)],
    bullet("Enforcement is guesswork.",
           "Teams patrol blindly. There is no map showing which spots"),
    [("", 17, WHITE, False), ("actually hurt traffic the most.", 17, WHITE, False)],
    bullet("Everything looks equally bad.",
           "Officers cannot tell which zone to fix first, so the worst"),
    [("", 17, WHITE, False), ("jams keep coming back.", 17, WHITE, False)],
], space_after=10, line_spacing=1.0)

# ---------------------------------------------------------------- Slide 3: Insight
s = slide()
kicker(s, "The insight")
title(s, "Not every wrong-parked car matters the same")
text(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.2),
     [[("A car blocking a busy main road at 6 PM is a big problem.", 20, WHITE, False)],
      [("The same car on a quiet street at midnight is not.", 20, WHITE, False)]],
     space_after=10)
box(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(1.5), fill=PANEL)
text(s, Inches(1.05), Inches(4.35), Inches(11.2), Inches(1.1),
     [[("So we stopped counting cars. ", 22, YELLOW, True),
       ("We measure how much each spot is actually hurting traffic flow.", 22, WHITE, True)]],
     anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- Slide 4: Solution (screenshot)
s = slide()
kicker(s, "Our solution")
title(s, "One live dashboard for the whole city", size=30)
h = pic_framed(s, os.path.join(SHOTS, "02-full-dashboard.png"), Inches(0.7), Inches(1.7), Inches(8.2))
text(s, Inches(9.2), Inches(1.7), Inches(3.5), Inches(5.0), [
    [("What you see", 16, YELLOW, True)],
    [("Hotspot map", 15, WHITE, True)],
    [("Live red / amber / green zones.", 13, MUTE, False)],
    [("Impact score", 15, WHITE, True)],
    [("0–100 score per spot.", 13, MUTE, False)],
    [("Priority queue", 15, WHITE, True)],
    [("Zones ranked by urgency.", 13, MUTE, False)],
    [("Situation report", 15, WHITE, True)],
    [("Plain-English brief per zone.", 13, MUTE, False)],
], space_after=8, line_spacing=1.0)

# ---------------------------------------------------------------- Slide 5: CIS
s = slide()
kicker(s, "How it works")
title(s, "The Congestion Impact Score (CIS)")
text(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.6),
     [[("A single 0–100 score that says how badly each spot hurts traffic. It is built from five things:",
        16, MUTE, False)]])
rows = [
    ("How many cars are wrongly parked", "42%", YELLOW),
    ("Time of day (rush hour vs night)", "22%", YELLOW),
    ("How important the road is", "18%", YELLOW),
    ("Risk of the jam spilling into junctions", "13%", YELLOW),
    ("History of repeat offenders", "5%", YELLOW),
]
y = Inches(2.7)
for i, (lbl, pct, col) in enumerate(rows):
    rb = box(s, Inches(0.7), y, Inches(8.2), Inches(0.62), fill=PANEL)
    text(s, Inches(0.95), y, Inches(6.6), Inches(0.62), [[(lbl, 15, WHITE, False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.4), y, Inches(1.3), Inches(0.62), [[(pct, 16, col, True)]],
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    y = y + Inches(0.72)
box(s, Inches(9.3), Inches(2.7), Inches(3.3), Inches(3.32), fill=None, line=BLUE, line_w=1.2)
text(s, Inches(9.55), Inches(2.95), Inches(2.85), Inches(2.9), [
    [("Why it's better", 15, YELLOW, True)],
    [("Every number is shown in", 14, WHITE, False)],
    [("the app — no black box.", 14, WHITE, False)],
    [("", 8, WHITE, False)],
    [("Score + how easy it is to", 14, WHITE, False)],
    [("reach the spot decides the", 14, WHITE, False)],
    [("Red / Amber / Green tier.", 14, WHITE, False)],
], space_after=4, line_spacing=1.0)

# ---------------------------------------------------------------- Slide 6: Transparency (screenshot)
s = slide()
kicker(s, "Fully explainable")
title(s, "No black box — every score is shown", size=30)
pic_framed(s, os.path.join(SHOTS, "03-transparency-raw-payload.png"), Inches(0.7), Inches(1.8), Inches(7.4))
text(s, Inches(8.5), Inches(2.2), Inches(4.2), Inches(4.0), [
    [("Click ", 17, WHITE, False), ("\u201cShow raw payload\u201d", 17, YELLOW, True),
     (" and you see the full breakdown behind every score.", 17, WHITE, False)],
    [("", 10, WHITE, False)],
    [("Judges, officers and engineers can all check exactly how a number was reached.",
      17, MUTE, False)],
], space_after=8)

# ---------------------------------------------------------------- Slide 7: Reacts (two screenshots)
s = slide()
kicker(s, "It reacts to reality")
title(s, "Busy at 6 PM, calm at midnight", size=30)
pic_framed(s, os.path.join(SHOTS, "01-dashboard-peak-1800.png"), Inches(0.7), Inches(1.95), Inches(5.8))
pic_framed(s, os.path.join(SHOTS, "04-offpeak-green-zones.png"), Inches(6.75), Inches(1.95), Inches(5.8))
text(s, Inches(0.7), Inches(5.55), Inches(5.8), Inches(0.5),
     [[("6 PM peak — ", 15, RED, True), ("red zones light up.", 15, WHITE, False)]])
text(s, Inches(6.75), Inches(5.55), Inches(5.8), Inches(0.5),
     [[("Midnight — ", 15, GREEN, True), ("most zones turn green.", 15, WHITE, False)]])
text(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.6),
     [[("Drag the time slider to any hour. The model is not \u201calways red\u201d — it follows real demand.",
        15, MUTE, False)]])

# ---------------------------------------------------------------- Slide 8: Action loop
s = slide()
kicker(s, "From alert to action")
title(s, "One click closes the loop")
steps = [
    ("1", "See", "Map and queue show the worst zones right now."),
    ("2", "Decide", "Read the plain-English situation report for that zone."),
    ("3", "Dispatch", "One click sends a patrol and estimates commuter-minutes saved."),
]
x = Inches(0.7)
for num, head, body in steps:
    box(s, x, Inches(2.4), Inches(3.75), Inches(2.4), fill=PANEL)
    text(s, x + Inches(0.3), Inches(2.65), Inches(1.0), Inches(0.8),
         [[(num, 34, YELLOW, True)]])
    text(s, x + Inches(0.3), Inches(3.45), Inches(3.2), Inches(0.5),
         [[(head, 20, WHITE, True)]])
    text(s, x + Inches(0.3), Inches(3.95), Inches(3.2), Inches(0.8),
         [[(body, 14, MUTE, False)]])
    x = x + Inches(4.0)
text(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.6),
     [[("Plus: export the queue to CSV for field teams, and a live feed updates the control room in real time.",
        15, MUTE, False)]])

# ---------------------------------------------------------------- Slide 9: Tech
s = slide()
kicker(s, "How it's built", color=BLUE)
title(s, "Simple, fast, reproducible")
pipe = ["Bengaluru data", "CIS engine", "API (REST + live feed)", "Dashboard"]
x = Inches(0.7)
for i, label in enumerate(pipe):
    box(s, x, Inches(2.3), Inches(2.55), Inches(0.95), fill=PANEL, line=BLUE, line_w=1.2)
    text(s, x, Inches(2.3), Inches(2.55), Inches(0.95), [[(label, 14, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(pipe) - 1:
        text(s, x + Inches(2.55), Inches(2.3), Inches(0.5), Inches(0.95),
             [[("\u2192", 22, YELLOW, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x = x + Inches(3.05)
text(s, Inches(0.7), Inches(3.9), Inches(11.9), Inches(2.5), [
    bullet("Backend:", "Python (FastAPI) computes the score and serves a live feed over WebSockets."),
    bullet("Frontend:", "React + Vite with a Leaflet map and live heatmap."),
    bullet("Data:", "Seeded synthetic dataset modeling real Bengaluru hotspots — same result every run."),
    bullet("AI:", "Briefings come from a clear rule engine, so results are explainable and repeatable."),
], space_after=12, line_spacing=1.0)

# ---------------------------------------------------------------- Slide 10: Impact + roadmap
s = slide()
kicker(s, "Impact & what's next")
title(s, "Smarter enforcement, clear next steps")
box(s, Inches(0.7), Inches(2.2), Inches(5.7), Inches(3.6), fill=PANEL)
text(s, Inches(1.0), Inches(2.45), Inches(5.1), Inches(3.2), [
    [("Impact today", 18, YELLOW, True)],
    [("\u2022  Fix the worst zones first.", 15, WHITE, False)],
    [("\u2022  Measure minutes saved for commuters.", 15, WHITE, False)],
    [("\u2022  Every decision is explainable.", 15, WHITE, False)],
], space_after=10)
box(s, Inches(6.9), Inches(2.2), Inches(5.7), Inches(3.6), fill=PANEL)
text(s, Inches(7.2), Inches(2.45), Inches(5.1), Inches(3.2), [
    [("What's next", 18, BLUE, True)],
    [("\u2022  Plug in real camera / ANPR feeds.", 15, WHITE, False)],
    [("\u2022  Mobile app for field officers.", 15, WHITE, False)],
    [("\u2022  Link to live traffic-signal data.", 15, WHITE, False)],
], space_after=10)

# ---------------------------------------------------------------- Slide 11: Thank you
s = slide()
box(s, 0, Inches(7.34), W, Inches(0.16), fill=YELLOW)
text(s, Inches(0.7), Inches(2.6), Inches(12), Inches(1.0), [[("Thank you", 48, WHITE, True)]])
text(s, Inches(0.74), Inches(3.7), Inches(12), Inches(0.6),
     [[("ParkSense — turning parking data into targeted enforcement.", 20, YELLOW, False)]])
text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(1.5), [
    [("Live demo:  http://localhost:5173", 15, MUTE, False)],
    [("Code + README in the repository.", 15, MUTE, False)],
], space_after=8)

for i, sl in enumerate(prs.slides):
    if 0 < i < len(prs.slides) - 1:
        footer(sl, i + 1)

prs.save(OUT)
print("Saved:", OUT)
